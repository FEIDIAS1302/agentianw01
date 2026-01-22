import streamlit as st
from openai import OpenAI
import json
import zipfile
import io
import datetime
import requests
from PIL import Image, ImageOps
from PyPDF2 import PdfReader
from pptx import Presentation

# --- デザイン設定 ---
st.set_page_config(page_title="AGENTIA for NUWORKS", layout="wide", page_icon="◾️")

st.markdown("""
<style>
    /* 全体のフォントと背景 */
    .stApp { font-family: 'Helvetica Neue', Arial, sans-serif; background-color: #ffffff; color: #1a1a1a; }
    /* ヘッダー周り */
    h1, h2, h3 { font-weight: 700 !important; letter-spacing: -0.05em !important; color: #000000 !important; }
    h1 { font-size: 3rem !important; margin-bottom: 0.5rem !important; }
    
    /* 入力フォーム */
    .stTextInput input, .stSelectbox div[data-baseweb="select"] { border-radius: 8px !important; border: 1px solid #e0e0e0 !important; padding: 0.5rem !important; }
    
    /* ボタンのスタイル */
    .stButton button { background-color: #000000 !important; color: #ffffff !important; border-radius: 30px !important; font-weight: bold !important; border: none !important; padding: 0.6rem 2rem !important; transition: all 0.3s ease; }
    .stButton button:hover { background-color: #333333 !important; box-shadow: 0 4px 12px rgba(0,0,0,0.15); transform: translateY(-2px); }

    /* 画像スタイル */
    img { border-radius: 12px; box-shadow: 0 4px 6px rgba(0,0,0,0.05); background-color: #f5f5f5; object-fit: contain; }
    hr { border-color: #f0f0f0; margin: 3rem 0; }
</style>
""", unsafe_allow_html=True)

# --- 設定 ---
# APIキーの取得（エラーハンドリング付き）
try:
    api_key = st.secrets["OPENAI_API_KEY"]
except:
    api_key = "" 

# --- データ定義 ---
BACKGROUNDS = {
    "bg_01": {"name": "Blue Abstract", "url": "assets/bg_01.jpg"},
    "bg_02": {"name": "White Marble", "url": "assets/bg_02.jpg"},
    "bg_03": {"name": "Modern Office", "url": "assets/bg_03.jpg"},
    "bg_04": {"name": "Tech Grid", "url": "assets/bg_04.jpg"},
}

AVATARS = {
    "avatar_a": {"name": "Avatar 01", "url": "assets/avat_01.png"},
    "avatar_b": {"name": "Avatar 02", "url": "assets/avat_02.png"},
    "avatar_c": {"name": "Avatar 03", "url": "assets/avat_03.png"},
    "avatar_d": {"name": "Avatar 04", "url": "assets/avat_04.png"},
}

BGMS = {
    "bgm_01": {"name": "Trust & Corporate", "desc": "信頼感のある明るいサウンド", "path": "assets/bgm1.mp3"},
    "bgm_02": {"name": "Innovation Tech", "desc": "先進的なデジタルビート", "path": "assets/bgm2.mp3"},
    "bgm_03": {"name": "Morning", "desc": "落ち着いた楽曲", "path": "assets/bgm3.mp3"},
    "bgm_04": {"name": "Future", "desc": "エネルギッシュなBGM", "path": "assets/bgm4.mp3"},
}

# --- ユーティリティ関数 ---

def load_image_from_url_or_path(path_or_url):
    """画像読み込み関数"""
    try:
        if path_or_url.startswith("http"):
            response = requests.get(path_or_url, stream=True)
            return Image.open(response.raw).convert("RGBA")
        else:
            return Image.open(path_or_url).convert("RGBA")
    except:
        # 画像がない場合のダミー
        return Image.new("RGBA", (1920, 1080), (240, 240, 240, 255))

def create_preview(bg_key, avatar_key, logo_upload):
    """プレビュー生成関数"""
    # 背景
    bg_img = load_image_from_url_or_path(BACKGROUNDS[bg_key]['url'])
    bg_img = bg_img.resize((1920, 1080))

    # アバター
    avatar_img = load_image_from_url_or_path(AVATARS[avatar_key]['url'])
    avatar_ratio = avatar_img.width / avatar_img.height
    new_h = 900
    new_w = int(new_h * avatar_ratio)
    avatar_img = avatar_img.resize((new_w, new_h))
    
    # アバター配置（中央下）
    x_pos = (1920 - new_w) // 2
    y_pos = 1080 - new_h
    bg_img.paste(avatar_img, (x_pos, y_pos), avatar_img)

    # ロゴ配置（左上）
    if logo_upload:
        logo_img = Image.open(logo_upload).convert("RGBA")
        l_ratio = logo_img.width / logo_img.height
        l_h = 80
        l_w = int(l_h * l_ratio)
        logo_img = logo_img.resize((l_w, l_h))
        bg_img.paste(logo_img, (60, 60), logo_img)

    return bg_img

def extract_text(file):
    """PDF/PPTXからテキスト抽出"""
    text = ""
    try:
        if file.name.endswith(".pdf"):
            reader = PdfReader(file)
            for page in reader.pages: text += page.extract_text()
        elif file.name.endswith(".pptx"):
            prs = Presentation(file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
    except: pass
    return text

def create_order_zip(order_data, logo_file, doc_file):
    """ZIPファイルを作成してメモリ上のデータとして返す"""
    zip_buffer = io.BytesIO()
    
    with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zf:
        # 1. JSON (設定ファイル)
        json_str = json.dumps(order_data, indent=4, ensure_ascii=False)
        zf.writestr("order.json", json_str)
        
        # 2. ロゴ画像
        if logo_file:
            logo_file.seek(0)
            ext = logo_file.name.split('.')[-1]
            zf.writestr(f"logo.{ext}", logo_file.read())
            
        # 3. 資料ファイル
        if doc_file:
            doc_file.seek(0)
            zf.writestr(doc_file.name, doc_file.read())
            
    return zip_buffer.getvalue()

def generate_script(text):
    """OpenAIによる台本生成"""
    if not text or len(text) < 10:
        return "エラー: 資料から文字を読み取れませんでした。"

    try:
        # Secretsからキーを取得して初期化
        client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "あなたはプロの動画構成作家です。"},
                {"role": "user", "content": f"""
                以下の資料を元に、企業の魅力が伝わる1分程度の動画台本を作成してください。
                
                【条件】
                - 読み上げ時間は約1分（文字数300〜400文字程度）
                - 丁寧すぎず、親しみやすい語り口で
                - 構成：導入（課題提起）→解決策（自社サービス）→実績・信頼性→結び
                
                【資料テキスト】
                {text[:15000]} 
                """}
            ]
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"AI生成エラー: {str(e)}"

# --- メインレイアウト ---

st.title("AGENTIA for NUWORKS")
st.markdown("Create your corporate video in minutes.")

col_input, col_preview = st.columns([1, 1.2], gap="large")

# === 左カラム：入力エリア ===
with col_input:
    st.markdown("### 1. Basic Info")
    project_id = st.text_input("Project ID", placeholder="NW10001")
    company_name = st.text_input("Company Name", placeholder="NuWorks Inc.")
    
    st.markdown("### 2. Assets")
    logo_file = st.file_uploader("Company Logo (PNG)", type=["png"])

    st.markdown("### 3. Visual Style")
    
    # 背景選択
    st.caption("Select Background")
    bg_keys = list(BACKGROUNDS.keys())
    bg_cols = st.columns(4)
    for i, key in enumerate(bg_keys):
        with bg_cols[i]:
            img = load_image_from_url_or_path(BACKGROUNDS[key]['url'])
            # サムネイルを正方形にクロップ
            min_side = min(img.width, img.height)
            square_img = ImageOps.fit(img, (min_side, min_side), centering=(0.5, 0.5))
            st.image(square_img, use_container_width=True)
            st.caption(BACKGROUNDS[key]['name'])
    
    bg_choice = st.radio("Choose Background", bg_keys, format_func=lambda x: BACKGROUNDS[x]['name'], horizontal=True, label_visibility="collapsed")

    # アバター選択
    st.caption("Select Avatar")
    av_keys = list(AVATARS.keys())
    c1, c2, c3, c4 = st.columns(4)
    with c1: st.image(AVATARS['avatar_a']['url']); st.caption("A")
    with c2: st.image(AVATARS['avatar_b']['url']); st.caption("B")
    with c3: st.image(AVATARS['avatar_c']['url']); st.caption("C")
    with c4: st.image(AVATARS['avatar_d']['url']); st.caption("D")
    avatar_choice = st.radio("Choose Model", av_keys, format_func=lambda x: AVATARS[x]['name'], horizontal=True)

    st.markdown("### 4. Audio")
    bgm_choice = st.selectbox("Background Music", list(BGMS.keys()), format_func=lambda x: BGMS[x]['name'])
    st.caption(f"♪ {BGMS[bgm_choice]['desc']}")
    
    # 試聴用プレイヤー
    try:
        st.audio(BGMS[bgm_choice]['path'], format="audio/mp3")
    except:
        pass
    
    st.markdown("### 5. Document")
    doc_file = st.file_uploader("Upload Company Profile (PDF/PPTX)", type=["pdf", "pptx"])
    
    st.divider()
    
    # --- 生成実行ボタン ---
    generate_clicked = st.button("Generate Script & Package", type="primary")

    if generate_clicked:
        if doc_file and company_name and project_id:
            with st.spinner("Analyzing document & Generating script..."):
                # 1. テキスト抽出
                txt = extract_text(doc_file)
                # 2. AI生成
                script = generate_script(txt)
                
                # 3. 結果をSession Stateに保存（再描画しても消えないように）
                st.session_state['generated_script'] = script
                st.session_state['generation_done'] = True
        else:
            st.error("⚠️ Project ID, Company Name, and Document are required!")

# === 右カラム：プレビューと結果エリア ===
with col_preview:
    st.markdown("### Preview")
    
    with st.container():
        # リアルタイム合成プレビュー
        preview_img = create_preview(bg_choice, avatar_choice, logo_file)
        st.image(preview_img, caption="Composite Preview", use_container_width=True)
        
        # 設定内容の要約表示
        st.markdown(f"""
        <div style="background-color:#f9f9f9; padding:1.5rem; border-radius:10px; border:1px solid #eee;">
            <p style="margin:0; font-size:0.9rem; color:#888;">CONFIGURATION</p>
            <h4 style="margin:0.5rem 0;">{BACKGROUNDS[bg_choice]['name']} / {AVATARS[avatar_choice]['name']}</h4>
            <p style="margin:0; font-size:0.9rem; color:#666;">🎵 BGM: {BGMS[bgm_choice]['name']}</p>
        </div>
        """, unsafe_allow_html=True)

    # --- 生成完了後の表示 ---
    if st.session_state.get('generation_done'):
        st.divider()
        st.subheader("✅ Generated Result")
        
        # 生成された台本（編集可能）
        final_script = st.text_area("Review Script", st.session_state['generated_script'], height=300)
        
        # --- ZIP作成処理 ---
        order_data = {
            "project_id": project_id,
            "company_name": company_name,
            "date": datetime.datetime.now().strftime("%Y%m%d"),
            "background_id": bg_choice,
            "avatar_id": avatar_choice,
            "bgm_id": bgm_choice,
            "script": final_script  # 編集後の台本を採用
        }
        
        # ZIPをバイナリとして作成
        zip_bytes = create_order_zip(order_data, logo_file, doc_file)
        file_name = f"{project_id}_{company_name}_Order.zip"
        
        st.markdown("### 📥 Download & Submit")
        st.info("以下の2ステップで納品してください：")

        # 2つのボタンを横並びに配置
        col_dl, col_dbx = st.columns([1, 1], gap="medium")
        
        with col_dl:
            # ステップ1: ダウンロード
            st.download_button(
                label="1️⃣ ZIPを保存 (Download)",
                data=zip_bytes,
                file_name=file_name,
                mime="application/zip",
                type="primary",
                use_container_width=True
            )
            
        with col_dbx:
            # ステップ2: Dropboxを開く
            st.link_button(
                label="2️⃣ Dropboxへアップロード 🚀",
                url="https://www.dropbox.com/request/DEvU9bqL8ncJP2H0kKzX",
                type="secondary",
                use_container_width=True
            )